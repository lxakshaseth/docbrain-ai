import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # Set 16:9 widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_slide_layout = prs.slide_layouts[6]  # Blank layout

    # Color Palette
    COLOR_BG = RGBColor(15, 23, 42)       # Dark Navy #0F172A
    COLOR_CARD = RGBColor(30, 41, 59)     # Slate #1E293B
    COLOR_CARD_BORDER = RGBColor(51, 65, 85) # #334155
    COLOR_CYAN = RGBColor(6, 182, 212)    # Cyan #06B6D4
    COLOR_LIGHT_CYAN = RGBColor(103, 232, 249) # #67E8F9
    COLOR_EMERALD = RGBColor(16, 185, 129)# Emerald #10B981
    COLOR_PURPLE = RGBColor(168, 85, 247) # Purple #A855F7
    COLOR_WHITE = RGBColor(248, 250, 252) # White #F8FAFC
    COLOR_MUTED = RGBColor(148, 163, 184) # Muted #94A3B8
    COLOR_DARK_TEXT = RGBColor(15, 23, 42)

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    def add_header(slide, title_text, category_text="DOCBRAIN AI • PROJECT PRESENTATION"):
        # Category Tag / Badge
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_CYAN
        p_cat.font.name = "Arial"

        # Main Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_WHITE
        p_title.font.name = "Arial"

    def add_card(slide, left, top, width, height, bg_color=COLOR_CARD, border_color=COLOR_CARD_BORDER):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background()
        return shape

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1)

    # Decorative hero card behind title
    add_card(slide1, 0.8, 1.2, 11.733, 5.2, bg_color=COLOR_CARD, border_color=COLOR_CYAN)

    # Hero Title Box
    t_box = slide1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.9), Inches(1.5))
    tf = t_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "DocBrain AI"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLOR_LIGHT_CYAN
    p.font.name = "Arial"

    p2 = tf.add_paragraph()
    p2.text = "Enterprise PDF Knowledge Base & Async RAG Architecture"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_WHITE
    p2.font.name = "Arial"

    # Subtitle / Description
    desc_box = slide1.shapes.add_textbox(Inches(1.2), Inches(3.4), Inches(10.9), Inches(1.2))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    p_desc = tf_desc.paragraphs[0]
    p_desc.text = "A production-grade, clean-architecture AI platform featuring Next.js 15, Node.js Express, Python FastAPI, LangGraph StateGraph, ChromaDB Vector Engine, and Redis Event Broker."
    p_desc.font.size = Pt(15)
    p_desc.font.color.rgb = COLOR_MUTED
    p_desc.font.name = "Arial"

    # Quick Metadata Badges
    meta_box = slide1.shapes.add_textbox(Inches(1.2), Inches(5.2), Inches(10.9), Inches(0.8))
    tf_meta = meta_box.text_frame
    tf_meta.word_wrap = True
    p_meta = tf_meta.paragraphs[0]
    p_meta.text = "⚡ Stack: Next.js 15 | Express | FastAPI | LangGraph | ChromaDB | Redis Pub/Sub | MongoDB Atlas"
    p_meta.font.size = Pt(13)
    p_meta.font.bold = True
    p_meta.font.color.rgb = COLOR_EMERALD
    p_meta.font.name = "Arial"

    # -------------------------------------------------------------
    # SLIDE 2: Executive Summary & Value Proposition
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide2)
    add_header(slide2, "Executive Summary & Core Value Proposition")

    # Card 1: Problem Statement
    add_card(slide2, 0.8, 1.6, 5.6, 5.2)
    tb = slide2.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🎯 The Challenge"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    items_prob = [
        "Traditional search fails on unstructured PDF documentation.",
        "Monolithic AI backend services create heavy memory bottlenecks and slow HTTP responses.",
        "Standard vector RAG often suffers from context hallucinations without keyword precision.",
        "Need for seamless, enterprise-grade multi-tenant authorization and document security."
    ]
    for item in items_prob:
        p_item = tf.add_paragraph()
        p_item.text = "• " + item
        p_item.font.size = Pt(13)
        p_item.font.color.rgb = COLOR_WHITE
        p_item.space_before = Pt(10)

    # Card 2: The Solution
    add_card(slide2, 6.9, 1.6, 5.6, 5.2)
    tb = slide2.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.2), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🚀 The DocBrain Solution"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_EMERALD

    items_sol = [
        "Decoupled Microservice Architecture: Asynchronous event broker eliminates backend HTTP blocking.",
        "Hybrid RAG Engine: Fuses Dense ChromaDB Vectors + Sparse BM25 Keywords via Reciprocal Rank Fusion.",
        "LangGraph StateGraph Workflow: Multi-step AI pipeline with self-correcting grounded answers & auto-generated follow-ups.",
        "Full Study Suite: Automatic summary generation, mind maps, flashcards, and gTTS audio overviews."
    ]
    for item in items_sol:
        p_item = tf.add_paragraph()
        p_item.text = "• " + item
        p_item.font.size = Pt(13)
        p_item.font.color.rgb = COLOR_WHITE
        p_item.space_before = Pt(10)

    # -------------------------------------------------------------
    # SLIDE 3: System Architecture Overview
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3)
    add_header(slide3, "System Architecture & Microservices Decoupling")

    # 4 Architecture Columns
    col_width = 2.7
    col_gap = 0.3
    left_start = 0.8
    top_pos = 1.6
    height_pos = 5.2

    services = [
        {
            "title": "1. Next.js 15 Web",
            "subtitle": "App Router & UI",
            "color": COLOR_CYAN,
            "points": ["React 18 & Zustand", "TanStack React Query", "Tailwind CSS & Shadcn", "Streaming Chat & Markdown", "Responsive UX Deck"]
        },
        {
            "title": "2. Node Express API",
            "subtitle": "API Gateway & Auth",
            "color": COLOR_LIGHT_CYAN,
            "points": ["Clean Architecture", "JWT & RBAC Security", "Zod DTO Validation", "Multer PDF Ingestion", "Swagger UI Docs"]
        },
        {
            "title": "3. Redis Pub/Sub",
            "subtitle": "Async Event Bus",
            "color": COLOR_PURPLE,
            "points": ["Zero Direct HTTP Coupling", "Channels: pdf:ingest", "Channels: rag:query", "Channels: rag:stream", "Resilient Task Queues"]
        },
        {
            "title": "4. Python AI Service",
            "subtitle": "FastAPI & RAG Engine",
            "color": COLOR_EMERALD,
            "points": ["LangGraph State Machine", "ChromaDB Vector Store", "Sparse BM25 Search", "Groq LPU & Gemini LLMs", "gTTS Audio & Study Tools"]
        }
    ]

    for idx, s in enumerate(services):
        c_left = left_start + idx * (col_width + col_gap)
        add_card(slide3, c_left, top_pos, col_width, height_pos)

        tb = slide3.shapes.add_textbox(Inches(c_left + 0.15), Inches(top_pos + 0.2), Inches(col_width - 0.3), Inches(height_pos - 0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_t = tf.paragraphs[0]
        p_t.text = s["title"]
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = s["color"]

        p_st = tf.add_paragraph()
        p_st.text = s["subtitle"]
        p_st.font.size = Pt(11)
        p_st.font.color.rgb = COLOR_MUTED

        for pt in s["points"]:
            p_pt = tf.add_paragraph()
            p_pt.text = "▪ " + pt
            p_pt.font.size = Pt(12)
            p_pt.font.color.rgb = COLOR_WHITE
            p_pt.space_before = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 4: Mandatory Redis Event Broker (Event-Driven Architecture)
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4)
    add_header(slide4, "Mandatory Redis Event Broker Architecture")

    # Card Top: Mandatory Architectural Rule
    add_card(slide4, 0.8, 1.5, 11.733, 1.2, bg_color=COLOR_CARD, border_color=COLOR_PURPLE)
    tb_rule = slide4.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(11.3), Inches(1.0))
    tf_rule = tb_rule.text_frame
    tf_rule.word_wrap = True
    p_r = tf_rule.paragraphs[0]
    p_r.text = "⚡ Architectural Mandate: Zero Synchronous Coupling"
    p_r.font.size = Pt(18)
    p_r.font.bold = True
    p_r.font.color.rgb = COLOR_PURPLE

    p_r2 = tf_rule.add_paragraph()
    p_r2.text = "Direct HTTP calls between Node.js API gateway and Python AI processing for heavy PDF ingestion and RAG streaming are strictly forbidden. All heavy workloads are dispatched asynchronously via Redis Pub/Sub channels."
    p_r2.font.size = Pt(13)
    p_r2.font.color.rgb = COLOR_WHITE
    p_r2.space_before = Pt(4)

    # Card Left: Event Workflow Channels
    add_card(slide4, 0.8, 2.9, 5.7, 4.0)
    tb_chan = slide4.shapes.add_textbox(Inches(1.0), Inches(3.1), Inches(5.3), Inches(3.6))
    tf_chan = tb_chan.text_frame
    tf_chan.word_wrap = True
    p_c = tf_chan.paragraphs[0]
    p_c.text = "📡 Redis Channels & Payload Flow"
    p_c.font.size = Pt(18)
    p_c.font.bold = True
    p_c.font.color.rgb = COLOR_CYAN

    channels = [
        ("pdf:ingest:request", "Node backend publishes PDF metadata & file path on upload."),
        ("pdf:ingest:response", "Python service processes PDF, populates vector store, and emits status."),
        ("rag:query:request", "Node publishes user query prompt & session ID."),
        ("rag:stream:chunk", "Python streams LLM response tokens back to Node for real-time SSE frontend delivery.")
    ]
    for ch, desc in channels:
        p_ch = tf_chan.add_paragraph()
        p_ch.text = f"• {ch}"
        p_ch.font.size = Pt(13)
        p_ch.font.bold = True
        p_ch.font.color.rgb = COLOR_LIGHT_CYAN
        p_ch.space_before = Pt(8)

        p_cd = tf_chan.add_paragraph()
        p_cd.text = desc
        p_cd.font.size = Pt(11)
        p_cd.font.color.rgb = COLOR_MUTED

    # Card Right: System Benefits
    add_card(slide4, 6.833, 2.9, 5.7, 4.0)
    tb_ben = slide4.shapes.add_textbox(Inches(7.033), Inches(3.1), Inches(5.3), Inches(3.6))
    tf_ben = tb_ben.text_frame
    tf_ben.word_wrap = True
    p_b = tf_ben.paragraphs[0]
    p_b.text = "🛡️ Key Architectural Advantages"
    p_b.font.size = Pt(18)
    p_b.font.bold = True
    p_b.font.color.rgb = COLOR_EMERALD

    bens = [
        "Non-blocking Node.js Event Loop: API remains instantly responsive under heavy PDF parsing loads.",
        "Independent Scaling: Python worker instances can scale horizontally according to vector embedding demand.",
        "Fault Isolation: Crashes in heavy AI LLM processing do not affect backend auth or session databases.",
        "Real-time Streaming Support: Seamless token-by-token pub/sub streaming directly to client WebSockets/SSE."
    ]
    for b in bens:
        p_bi = tf_ben.add_paragraph()
        p_bi.text = "✔ " + b
        p_bi.font.size = Pt(12)
        p_bi.font.color.rgb = COLOR_WHITE
        p_bi.space_before = Pt(10)

    # -------------------------------------------------------------
    # SLIDE 5: Python AI Engine & LangGraph StateGraph Workflow
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide5)
    add_header(slide5, "LangGraph StateGraph Workflow (AI Engine Pipeline)")

    # 3 Workflow Nodes
    node_w = 3.64
    node_gap = 0.4
    n_top = 1.6
    n_height = 5.2

    nodes = [
        {
            "step": "STEP 1",
            "title": "hybrid_retrieve",
            "subtitle": "Multi-Engine Context Retrieval",
            "color": COLOR_CYAN,
            "details": [
                "Executes ChromaDB Dense Vector embedding query.",
                "Executes BM25 Sparse Keyword search across document chunks.",
                "Merges rankings using Reciprocal Rank Fusion (RRF).",
                "Extracts top-k most relevant grounded context snippets."
            ]
        },
        {
            "step": "STEP 2",
            "title": "generate_answer",
            "subtitle": "Anti-Hallucination LLM Generation",
            "color": COLOR_LIGHT_CYAN,
            "details": [
                "Injects strict anti-hallucination system prompt templates.",
                "Fetches chat memory history for session state.",
                "Invokes Groq (Llama-3.1-70b) or Google Gemini 1.5 Flash.",
                "Enforces strict attribution to uploaded PDF context."
            ]
        },
        {
            "step": "STEP 3",
            "title": "generate_followups",
            "subtitle": "Grounded Question Suggestion",
            "color": COLOR_EMERALD,
            "details": [
                "Parses generated response and document context.",
                "Generates 3 logical, contextually grounded follow-up prompts.",
                "Returns structured JSON response payload to client.",
                "Enhances user engagement and exploration."
            ]
        }
    ]

    for idx, nd in enumerate(nodes):
        n_left = 0.8 + idx * (node_w + node_gap)
        add_card(slide5, n_left, n_top, node_w, n_height)

        tb = slide5.shapes.add_textbox(Inches(n_left + 0.2), Inches(n_top + 0.2), Inches(node_w - 0.4), Inches(n_height - 0.4))
        tf = tb.text_frame
        tf.word_wrap = True

        p_s = tf.paragraphs[0]
        p_s.text = nd["step"]
        p_s.font.size = Pt(11)
        p_s.font.bold = True
        p_s.font.color.rgb = nd["color"]

        p_t = tf.add_paragraph()
        p_t.text = nd["title"]
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_WHITE

        p_sub = tf.add_paragraph()
        p_sub.text = nd["subtitle"]
        p_sub.font.size = Pt(12)
        p_sub.font.color.rgb = COLOR_MUTED

        for d in nd["details"]:
            p_d = tf.add_paragraph()
            p_d.text = "• " + d
            p_d.font.size = Pt(12)
            p_d.font.color.rgb = COLOR_WHITE
            p_d.space_before = Pt(10)

    # -------------------------------------------------------------
    # SLIDE 6: Hybrid Retrieval Engine (BM25 + ChromaDB Vector RRF)
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide6)
    add_header(slide6, "Hybrid Retrieval Strategy: Vector + Keyword RRF")

    # Card Left: Dense vs Sparse
    add_card(slide6, 0.8, 1.6, 5.6, 5.2)
    tb_l = slide6.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True

    p = tf_l.paragraphs[0]
    p.text = "🧠 Dual-Retrieval Mechanics"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    mechanics = [
        ("ChromaDB Dense Embeddings", "Captures high-level semantic meaning, synonyms, and conceptual intent using SentenceTransformers / Gemini embeddings."),
        ("Sparse BM25 Keyword Search", "Captures exact alphanumeric matches, technical jargon, proper nouns, and specific formula codes."),
        ("Reciprocal Rank Fusion (RRF)", "RRF Score = ∑ 1 / (k + rank_i). Fuses rank order lists from both search engines without needing raw score normalization.")
    ]
    for title, desc in mechanics:
        p_t = tf_l.add_paragraph()
        p_t.text = "🔹 " + title
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_WHITE
        p_t.space_before = Pt(10)

        p_d = tf_l.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = COLOR_MUTED

    # Card Right: Anti-Hallucination & Quality Metrics
    add_card(slide6, 6.9, 1.6, 5.6, 5.2)
    tb_r = slide6.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.2), Inches(4.8))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True

    p_r = tf_r.paragraphs[0]
    p_r.text = "🎯 Accuracy & Anti-Hallucination"
    p_r.font.size = Pt(20)
    p_r.font.bold = True
    p_r.font.color.rgb = COLOR_EMERALD

    metrics = [
        "Zero-Context Fallback: System explicitly rejects queries outside PDF knowledge boundaries.",
        "Citation Metadata: Every response includes exact page numbers and PDF chunk references.",
        "Chunking Strategy: Recursive character text splitter (1000 chunk size, 200 overlap) preserves contextual continuity.",
        "Multi-LLM Resilience: Dynamic failover between Groq LPU (Llama-3.1-70b) and Google Gemini 1.5 Flash."
    ]
    for m in metrics:
        p_m = tf_r.add_paragraph()
        p_m.text = "✔ " + m
        p_m.font.size = Pt(13)
        p_m.font.color.rgb = COLOR_WHITE
        p_m.space_before = Pt(12)

    # -------------------------------------------------------------
    # SLIDE 7: Advanced Features & AI Study Suite
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide7)
    add_header(slide7, "Advanced Intelligence & AI Study Suite")

    # 4 Quadrant Cards
    q_w = 5.6
    q_h = 2.45
    coords = [
        (0.8, 1.6, "📄 Automated Document Summarization", COLOR_CYAN, [
            "Generates executive summaries, key takeaways, and section breakdowns upon PDF upload.",
            "Powered by specialized summary templates in Python AI service."
        ]),
        (6.9, 1.6, "🗺️ Mind Map & Flashcard Generator", COLOR_EMERALD, [
            "Extracts key entity relationships and generates interactive visual mind maps.",
            "Automatically creates interactive study flashcards for quick revision."
        ]),
        (0.8, 4.35, "📊 Multi-PDF Comparison Engine", COLOR_LIGHT_CYAN, [
            "Side-by-side comparison between multiple uploaded PDF documents.",
            "Highlights key differences, contradictory claims, and shared topics."
        ]),
        (6.9, 4.35, "🔊 Audio Overview & Public Share Links", COLOR_PURPLE, [
            "Text-to-speech audio overview generated via gTTS engine.",
            "Public read-only sharing via secure tokenized URLs (/share/:token)."
        ])
    ]

    for left, top, title, color, bullet_list in coords:
        add_card(slide7, left, top, q_w, q_h)
        tb = slide7.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.15), Inches(q_w - 0.3), Inches(q_h - 0.3))
        tf = tb.text_frame
        tf.word_wrap = True

        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = color

        for b in bullet_list:
            p_b = tf.add_paragraph()
            p_b.text = "• " + b
            p_b.font.size = Pt(11)
            p_b.font.color.rgb = COLOR_WHITE
            p_b.space_before = Pt(4)

    # -------------------------------------------------------------
    # SLIDE 8: Complete Technology Stack Matrix
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide8)
    add_header(slide8, "Comprehensive Technology Stack Breakdown")

    # Table Layout
    rows = 6
    cols = 3
    t_left = Inches(0.8)
    t_top = Inches(1.6)
    t_width = Inches(11.733)
    t_height = Inches(5.2)

    table_shape = slide8.shapes.add_table(rows, cols, t_left, t_top, t_width, t_height)
    table = table_shape.table

    # Column widths
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(3.5)
    table.columns[2].width = Inches(5.733)

    table_data = [
        ["Layer", "Technologies Used", "Key Responsibilities & Role"],
        ["Frontend UI", "Next.js 15 (App Router), React 18, Tailwind CSS, shadcn/ui", "Responsive UI, Server Components, Zustand state, TanStack Query data fetching"],
        ["Backend API Gateway", "Node.js, Express, TypeScript, Zod, Multer, Swagger UI", "Authentication, RBAC, DTO validation, document metadata, Swagger OpenAPI docs"],
        ["AI Microservice", "Python 3.11, FastAPI, LangChain, LangGraph, PyPDF, gTTS", "RAG pipeline, state graph orchestration, document parsing, audio overview creation"],
        ["Vector & Document Store", "ChromaDB (Vector Store), MongoDB Atlas (Cloud NoSQL)", "Persistent high-dimensional embeddings, user accounts, document metadata, chat history"],
        ["Event Bus & Infra", "Redis Pub/Sub, Docker, Docker Compose, Upstash Redis", "Asynchronous messaging, decoupling API from AI processing, containerization"]
    ]

    for r_idx, row in enumerate(table_data):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid()
            if r_idx == 0:
                cell.fill.fore_color.rgb = COLOR_CARD_BORDER
            else:
                cell.fill.fore_color.rgb = COLOR_CARD

            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.name = "Arial"
            if r_idx == 0:
                p.font.bold = True
                p.font.size = Pt(13)
                p.font.color.rgb = COLOR_CYAN
            else:
                p.font.size = Pt(11)
                if c_idx == 0:
                    p.font.bold = True
                    p.font.color.rgb = COLOR_LIGHT_CYAN
                else:
                    p.font.color.rgb = COLOR_WHITE

    # -------------------------------------------------------------
    # SLIDE 9: Security, Auth & Data Integrity
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide9)
    add_header(slide9, "Security, Authentication & Data Integrity")

    # 3 Horizontal Cards
    card_h = 1.55
    top_base = 1.6
    gap_base = 0.2

    sec_cards = [
        ("🔐 Multi-Tenant JWT Authentication & RBAC", COLOR_CYAN, [
            "Stateless JSON Web Tokens (JWT) with password hashing via bcrypt.",
            "Role-Based Access Control (RBAC) middleware enforcing document access boundaries per user."
        ]),
        ("🛡️ Strict Input Validation & DTO Hygiene", COLOR_EMERALD, [
            "Zod validation schemas on all incoming HTTP requests in Node.js backend.",
            "Pydantic schemas enforcing strict event payloads across Redis Pub/Sub channels."
        ]),
        ("📁 Secure File Processing & Storage Isolation", COLOR_LIGHT_CYAN, [
            "Multer upload middleware sanitizes filename extensions and restricts MIME types.",
            "Vector database namespaces isolate user document embeddings from cross-tenant access."
        ])
    ]

    for idx, (title, color, bullets) in enumerate(sec_cards):
        c_top = top_base + idx * (card_h + gap_base)
        add_card(slide9, 0.8, c_top, 11.733, card_h)

        tb = slide9.shapes.add_textbox(Inches(1.0), Inches(c_top + 0.15), Inches(11.333), Inches(card_h - 0.3))
        tf = tb.text_frame
        tf.word_wrap = True

        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(17)
        p_t.font.bold = True
        p_t.font.color.rgb = color

        for b in bullets:
            p_b = tf.add_paragraph()
            p_b.text = "• " + b
            p_b.font.size = Pt(12)
            p_b.font.color.rgb = COLOR_WHITE
            p_b.space_before = Pt(4)

    # -------------------------------------------------------------
    # SLIDE 10: API Architecture & Health Monitoring
    # -------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide10)
    add_header(slide10, "API Architecture & System Health Checks")

    # Card Left: Swagger OpenAPI
    add_card(slide10, 0.8, 1.6, 5.6, 5.2)
    tb_sw = slide10.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.8))
    tf_sw = tb_sw.text_frame
    tf_sw.word_wrap = True

    p = tf_sw.paragraphs[0]
    p.text = "📖 Swagger UI OpenAPI Docs"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    sw_items = [
        "Interactive Documentation: Served directly at http://localhost:5000/api-docs.",
        "Complete Endpoint Coverage: Full documentation of Auth, Documents, and Chat endpoints.",
        "Try-It-Out Sandbox: Built-in authorization header configuration for seamless API testing."
    ]
    for item in sw_items:
        p_i = tf_sw.add_paragraph()
        p_i.text = "• " + item
        p_i.font.size = Pt(12)
        p_i.font.color.rgb = COLOR_WHITE
        p_i.space_before = Pt(10)

    # Card Right: Health Endpoints
    add_card(slide10, 6.9, 1.6, 5.6, 5.2)
    tb_hp = slide10.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.2), Inches(4.8))
    tf_hp = tb_hp.text_frame
    tf_hp.word_wrap = True

    p_h = tf_hp.paragraphs[0]
    p_h.text = "🧪 Multi-Tier Health Probe Endpoints"
    p_h.font.size = Pt(20)
    p_h.font.bold = True
    p_h.font.color.rgb = COLOR_EMERALD

    h_items = [
        ("GET /health", "Node.js Backend health check verifying server uptime and DB state."),
        ("GET /health/redis", "Verifies live connectivity to Redis Pub/Sub broker."),
        ("GET /api/v1/health", "Python FastAPI health check verifying ChromaDB and LLM key readiness.")
    ]
    for ep, desc in h_items:
        p_ep = tf_hp.add_paragraph()
        p_ep.text = "🟢 " + ep
        p_ep.font.size = Pt(14)
        p_ep.font.bold = True
        p_ep.font.color.rgb = COLOR_LIGHT_CYAN
        p_ep.space_before = Pt(10)

        p_ed = tf_hp.add_paragraph()
        p_ed.text = desc
        p_ed.font.size = Pt(11)
        p_ed.font.color.rgb = COLOR_MUTED

    # -------------------------------------------------------------
    # SLIDE 11: Production Cloud Deployment Strategy
    # -------------------------------------------------------------
    slide11 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide11)
    add_header(slide11, "Cloud Deployment & Infrastructure Blueprint")

    # 4 Cards for Hosting Targets
    card_width_d = 5.6
    card_height_d = 2.45

    deploy_targets = [
        (0.8, 1.6, "▲ Frontend Hosting: Vercel", COLOR_CYAN, [
            "Deployed via Vercel GitHub Integration (apps/frontend).",
            "Environment configuration: NEXT_PUBLIC_API_URL pointing to backend.",
            "Global CDN caching for static assets & SSG pages."
        ]),
        (6.9, 1.6, "🔷 Microservices: Render Blueprint", COLOR_EMERALD, [
            "Automated multi-service blueprint via render.yaml configuration.",
            "Deploys Node.js Express backend & Python FastAPI AI microservice.",
            "Auto-scaling web service instances with custom health checks."
        ]),
        (0.8, 4.35, "🍃 Database: MongoDB Atlas Cloud", COLOR_LIGHT_CYAN, [
            "Free M0 Cloud Cluster for persistent document & user metadata.",
            "Network access restricted with secure database user credentials."
        ]),
        (6.9, 4.35, "⚡ Cache & Event Bus: Upstash Redis", COLOR_PURPLE, [
            "Serverless Redis instance supporting Pub/Sub messaging.",
            "Low-latency cross-cloud event streaming between Render services."
        ])
    ]

    for left, top, title, color, bullets in deploy_targets:
        add_card(slide11, left, top, card_width_d, card_height_d)
        tb = slide11.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.15), Inches(card_width_d - 0.3), Inches(card_height_d - 0.3))
        tf = tb.text_frame
        tf.word_wrap = True

        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = color

        for b in bullets:
            p_b = tf.add_paragraph()
            p_b.text = "• " + b
            p_b.font.size = Pt(11)
            p_b.font.color.rgb = COLOR_WHITE
            p_b.space_before = Pt(4)

    # -------------------------------------------------------------
    # SLIDE 12: Project Summary & Key Achievements
    # -------------------------------------------------------------
    slide12 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide12)
    add_header(slide12, "Summary & Key Architectural Achievements")

    # Big Card Summary
    add_card(slide12, 0.8, 1.6, 11.733, 5.2, bg_color=COLOR_CARD, border_color=COLOR_CYAN)
    tb_sum = slide12.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(11.133), Inches(4.8))
    tf_sum = tb_sum.text_frame
    tf_sum.word_wrap = True

    p = tf_sum.paragraphs[0]
    p.text = "🏆 Key Milestones & Production Highlights"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_LIGHT_CYAN

    achievements = [
        "Production-Grade Decoupling: Fully event-driven microservices architecture via Redis Pub/Sub.",
        "Advanced Hybrid RAG: Reciprocal Rank Fusion of ChromaDB Dense Vector Search + BM25 Sparse Keyword Search.",
        "Resilient AI Workflows: LangGraph state machine with automatic follow-up generation and zero-hallucination prompts.",
        "Comprehensive Feature Suite: Full support for RAG chat, document summaries, interactive mind maps, flashcards, PDF comparison, and gTTS audio overviews.",
        "DevOps & Security Ready: Containerized via Docker Compose, automated Render Blueprint, JWT/RBAC security, and interactive Swagger OpenAPI documentation."
    ]

    for a in achievements:
        p_a = tf_sum.add_paragraph()
        p_a.text = "✅ " + a
        p_a.font.size = Pt(14)
        p_a.font.color.rgb = COLOR_WHITE
        p_a.space_before = Pt(12)

    output_path = r"c:\Users\lxaks\OneDrive\Desktop\intern\DocBrain_AI_Project_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_deck()
